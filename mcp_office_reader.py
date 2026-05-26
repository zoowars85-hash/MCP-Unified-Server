#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MCP Office Reader v3.1
Secure extraction and conversion for Word, Excel, PowerPoint documents.
"""
import os
import sys
import subprocess
from pathlib import Path
from typing import List, Dict, Optional, Any
from mcp_shared import (
    _log, normalize_path, _ensure_allowed,
    BaseMCPServer, conversation_memory, dialog_ctx
)

def read_docx(file_path: str, include_headers: bool = True) -> Dict:
    try:
        from docx import Document
    except ImportError:
        return {"error": "python-docx not installed. Install with: pip install python-docx"}

    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "read_docx")
    if not p.is_file():
        return {"error": f"File not found: {file_path}"}

    doc = Document(str(p))
    content = []
    for para in doc.paragraphs:
        content.append({
            "text": para.text,
            "style": para.style.name if include_headers else None,
            "alignment": para.alignment.name if para.alignment else None
        })

    tables = []
    for i, table in enumerate(doc.tables):
        table_data = [[cell.text for cell in row.cells] for row in table.rows]
        tables.append({"table_index": i, "data": table_data})

    result = {"path": str(p), "filename": p.name, "paragraphs": content, "tables": tables}
    conversation_memory.add(
        op="read_docx", paths={"file": str(p)}, status="success", dialog=dialog_ctx.get(),
        context=f"Extracted {len(content)} paragraphs and {len(tables)} tables from {p.name}"
    )
    return result

def read_excel(file_path: str, sheet_name: Optional[str] = None, max_rows: int = 1000) -> Dict:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return {"error": "openpyxl not installed. Install with: pip install openpyxl"}

    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "read_excel")
    if not p.is_file():
        return {"error": f"File not found: {file_path}"}

    wb = load_workbook(str(p), read_only=True, data_only=True)
    ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

    if ws is None:
        return {"error": f"Sheet '{sheet_name}' not found"}

    data = []
    headers = None
    row_count = 0
    for row in ws.iter_rows(values_only=True):
        row_count += 1
        if row_count > max_rows:
            break
        if headers is None:
            headers = list(row)
            continue
        row_dict = {h: v for h, v in zip(headers, row) if h is not None}
        data.append(row_dict)

    wb.close()
    result = {"path": str(p), "sheet": ws.title, "headers": headers, "rows_extracted": len(data), "data": data}
    conversation_memory.add(
        op="read_excel", paths={"file": str(p)}, status="success", dialog=dialog_ctx.get(),
        context=f"Extracted {len(data)} rows from sheet '{ws.title}'"
    )
    return result

def read_pptx(file_path: str, slide_numbers: Optional[List[int]] = None) -> Dict:
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed. Install with: pip install python-pptx"}

    p = Path(normalize_path(file_path))
    _ensure_allowed(p, "read_pptx")
    if not p.is_file():
        return {"error": f"File not found: {file_path}"}

    prs = Presentation(str(p))
    slides_data = []
    for i, slide in enumerate(prs.slides):
        slide_idx = i + 1
        if slide_numbers and slide_idx not in slide_numbers:
            continue

        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(paragraph.text)

        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        except Exception:
            pass

        slides_data.append({
            "slide_number": slide_idx,
            "text": "\n".join(texts),
            "notes": notes
        })

    result = {"path": str(p), "total_slides": len(prs.slides), "extracted_slides": slides_data}
    conversation_memory.add(
        op="read_pptx", paths={"file": str(p)}, status="success", dialog=dialog_ctx.get(),
        context=f"Extracted text from {len(slides_data)} slides in {p.name}"
    )
    return result

def export_to_pdf(input_path: str, output_path: str) -> Dict:
    src = Path(normalize_path(input_path))
    dst = Path(normalize_path(output_path))
    _ensure_allowed(src, "export_to_pdf")
    _ensure_allowed(dst.parent, "export_to_pdf")
    if not src.is_file():
        return {"error": f"Input file not found: {input_path}"}

    cmd = ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(dst.parent), str(src)]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=90)
        if result.returncode != 0:
            alt_cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(dst.parent), str(src)]
            result = subprocess.run(alt_cmd, capture_output=True, text=True, timeout=90)

        if result.returncode != 0:
            return {"error": f"LibreOffice conversion failed: {result.stderr.strip()}"}

        expected_out = dst.parent / f"{src.stem}.pdf"
        if expected_out.exists():
            if expected_out != dst:
                expected_out.rename(dst)
            conversation_memory.add(
                op="export_to_pdf", paths={"src": str(src), "dst": str(dst)},
                status="success", dialog=dialog_ctx.get(), context=f"Converted {src.name} to PDF"
            )
            return {"status": "success", "output": str(dst)}
        return {"error": "Conversion completed but output file not found."}
    except FileNotFoundError:
        return {"error": "LibreOffice/soffice not found. Install it to enable PDF export."}
    except Exception as e:
        return {"error": str(e)}

server = BaseMCPServer("office-reader", "3.1")
server.register_tool("read_docx", {
    "description": "Extract text, tables, and styles from DOCX files",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "include_headers": {"type": "boolean", "default": True}
        },
        "required": ["file_path"]
    }
}, lambda **kw: read_docx(kw["file_path"], kw.get("include_headers", True)))

server.register_tool("read_excel", {
    "description": "Extract sheet data from Excel files to JSON structure",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "sheet_name": {"type": "string"},
            "max_rows": {"type": "integer", "default": 1000}
        },
        "required": ["file_path"]
    }
}, lambda **kw: read_excel(kw["file_path"], kw.get("sheet_name"), kw.get("max_rows", 1000)))

server.register_tool("read_pptx", {
    "description": "Extract text and speaker notes from PowerPoint presentations",
    "inputSchema": {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "slide_numbers": {"type": "array", "items": {"type": "integer"}}
        },
        "required": ["file_path"]
    }
}, lambda **kw: read_pptx(kw["file_path"], kw.get("slide_numbers")))

server.register_tool("export_to_pdf", {
    "description": "Convert supported office documents to PDF via LibreOffice",
    "inputSchema": {
        "type": "object",
        "properties": {
            "input_path": {"type": "string"},
            "output_path": {"type": "string"}
        },
        "required": ["input_path", "output_path"]
    }
}, lambda **kw: export_to_pdf(kw["input_path"], kw["output_path"]))

if __name__ == "__main__":
    server.run()